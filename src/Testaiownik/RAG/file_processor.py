from typing import List
from docx import Document
import pdfplumber
from utils.logger import logger
from pptx import Presentation


def extract_text_from_pdf(pdf_path: str) -> List[tuple[str, int]]:
    """Extracts all text from a PDF file with page numbers."""
    with pdfplumber.open(pdf_path) as pdf:
        results = []
        for i, page in enumerate(pdf.pages, start=1):
            page_text = page.extract_text()
            if page_text:
                results.append((page_text, i))
        return results


def extract_text_from_txt(txt_path: str) -> str:
    """Extracts all text from a TXT file."""
    try:
        with open(txt_path, "r", encoding="utf-8") as file:
            text = file.read()
        return text
    except Exception as e:
        logger.error(f"Error reading the text file: {e}")
        return ""


def extract_text_from_pptx(pptx_path: str) -> List[tuple[str, int]]:
    """Extracts text from a PPTX file with slide numbers."""
    results = []
    try:
        presentation = Presentation(pptx_path)
        for i, slide in enumerate(presentation.slides, start=1):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            if slide_text.strip():
                results.append((slide_text.strip(), i))
        return results
    except Exception as e:
        logger.error(f"Błąd podczas wczytywania pliku PPTX: {e}")
        return []


def extract_text_from_docx(docx_path: str) -> str:
    """Extracts text from a DOCX file."""
    try:
        doc = Document(docx_path)

        text = "\n".join([para.text for para in doc.paragraphs])

        return text
    except Exception as e:
        logger.error(f"Błąd podczas wczytywania pliku DOCX: {e}")
        return ""
